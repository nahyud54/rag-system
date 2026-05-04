"""Stage 1: ETL - Extract text from PPTX files"""

import logging
from pathlib import Path
from typing import Any, Dict, List

from pptx import Presentation

logger = logging.getLogger(__name__)


class PPTXExtractor:
    """Extracts text content from PPTX files, one document per slide."""

    def extract_from_file(self, file_path: str) -> List[Dict[str, Any]]:
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        prs = Presentation(str(path))
        documents: List[Dict[str, Any]] = []
        course_name = path.stem

        for slide_idx, slide in enumerate(prs.slides, 1):
            texts: List[str] = []

            # Title first
            if slide.shapes.title and slide.shapes.title.text.strip():
                texts.append(slide.shapes.title.text.strip())

            # Text frames (skip title shape to avoid duplicate)
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            texts.append(line)

            # Table cells
            for shape in slide.shapes:
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            texts.append(" | ".join(cells))

            slide_text = "\n".join(texts).strip()
            if len(slide_text) < 20:
                continue

            documents.append({
                "text": slide_text,
                "metadata": {
                    "source": str(path),
                    "course_name": course_name,
                    "slide_number": slide_idx,
                    "total_slides": len(prs.slides),
                    "file_name": path.name,
                },
            })

        logger.info(f"Extracted {len(documents)} slides from {path.name}")
        return documents

    def extract_from_directory(self, directory: str) -> List[Dict[str, Any]]:
        dir_path = Path(directory)
        if not dir_path.exists():
            raise FileNotFoundError(f"Directory not found: {directory}")

        all_docs: List[Dict[str, Any]] = []
        pptx_files = sorted(dir_path.glob("*.pptx")) + sorted(dir_path.glob("*.ppt"))

        for fp in pptx_files:
            try:
                docs = self.extract_from_file(str(fp))
                all_docs.extend(docs)
            except Exception as e:
                logger.error(f"Failed to process {fp.name}: {e}")

        logger.info(f"Total slides extracted from directory: {len(all_docs)}")
        return all_docs
